from pptx import Presentation
from pathlib import Path
from app.ingestion.base import BaseParser
from app.models.document import FileFormat, ParsedDocument
from app.core.exceptions import ParseFailureError


class PPTParser(BaseParser):
    @property
    def supported_format(self) -> FileFormat:
        return FileFormat.PPT

    @property
    def supported_extensions(self) -> set[str]:
        return {"pptx"}

    def parse(self, file_path: Path) -> ParsedDocument:
        try:
            prs = Presentation(file_path)
            flat_texts = []
            structured_slides = []

            for i, slide in enumerate(prs.slides, start=1):
                title = slide.shapes.title.text if slide.shapes.title else ""

                full_text = "\n".join(
                    [shape.text for shape in slide.shapes if hasattr(shape, "text")]
                )

                flat_texts.append(full_text)
                structured_slides.append(
                    {"slide_number": i, "title": title, "content": full_text}
                )

            return ParsedDocument(
                file_name=file_path.name,
                format=self.supported_format,
                page_count=len(prs.slides),
                text_content=flat_texts,
                metadata={
                    "slide_count": len(prs.slides),
                    "structured_slides": structured_slides,
                },
            )
        except Exception as e:
            raise ParseFailureError(file_path.name, str(e))
