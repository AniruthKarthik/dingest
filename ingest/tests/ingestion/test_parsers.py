import csv
from pathlib import Path

import pandas as pd
import pytest
from docx import Document
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

from app.core.exceptions import ParseFailureError
from app.ingestion.csv_parser import CSVParser
from app.ingestion.excel_parser import ExcelParser
from app.ingestion.pdf_parser import PDFParser
from app.ingestion.ppt_parser import PPTParser
from app.ingestion.word_parser import WordParser
from app.models.document import FileFormat


@pytest.fixture
def pdf_file(tmp_path: Path) -> Path:
    path = tmp_path / "sample.pdf"
    c = canvas.Canvas(str(path), pagesize=letter)
    c.drawString(72, 720, "DocuHub PDF test page one.")
    c.showPage()
    c.drawString(72, 720, "Page two content.")
    c.showPage()
    c.save()
    return path


@pytest.fixture
def excel_file(tmp_path: Path) -> Path:
    path = tmp_path / "sample.xlsx"
    df = pd.DataFrame({"name": ["Alice", "Bob"], "score": [95, 87]})
    df.to_excel(path, index=False, sheet_name="Results")
    return path


@pytest.fixture
def csv_file(tmp_path: Path) -> Path:
    path = tmp_path / "sample.csv"
    with open(path, "w", newline="") as f:
        writer = csv.DictWriter(f, fieldnames=["city", "population"])
        writer.writeheader()
        writer.writerow({"city": "Chennai", "population": 7000000})
    return path


@pytest.fixture
def word_file(tmp_path: Path) -> Path:
    path = tmp_path / "sample.docx"
    doc = Document()
    doc.add_heading("Test Heading", level=1)
    doc.add_paragraph("This is a test paragraph.")
    doc.save(str(path))
    return path


@pytest.fixture
def ppt_file(tmp_path: Path) -> Path:
    path = tmp_path / "sample.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Test Slide Title"
    slide.placeholders[1].text = "Subtitle text"
    prs.save(str(path))
    return path


class TestPDFParser:
    def test_format_and_extensions(self):
        p = PDFParser()
        assert p.supported_format == FileFormat.PDF
        assert "pdf" in p.supported_extensions

    def test_can_handle(self):
        p = PDFParser()
        assert p.can_handle("pdf") is True
        assert p.can_handle(".pdf") is True
        assert p.can_handle("xlsx") is False

    def test_parse_returns_correct_structure(self, pdf_file):
        result = PDFParser().parse(pdf_file)
        assert result.file_name == "sample.pdf"
        assert result.format == FileFormat.PDF
        assert result.page_count == 2
        assert len(result.text_content) == 2
        assert "DocuHub" in result.text_content[0]

    def test_parse_invalid_file_raises(self, tmp_path):
        bad = tmp_path / "bad.pdf"
        bad.write_text("not a pdf")
        with pytest.raises(ParseFailureError):
            PDFParser().parse(bad)


class TestExcelParser:
    def test_format_and_extensions(self):
        p = ExcelParser()
        assert p.supported_format == FileFormat.EXCEL
        assert "xlsx" in p.supported_extensions

    def test_parse_returns_sheets(self, excel_file):
        result = ExcelParser().parse(excel_file)
        assert result.format == FileFormat.EXCEL
        assert result.sheets is not None
        assert "Results" in result.sheets
        rows = result.sheets["Results"]
        assert len(rows) == 2
        assert rows[0]["name"] == "Alice"
        assert rows[1]["score"] == 87

    def test_parse_page_count_equals_sheet_count(self, excel_file):
        result = ExcelParser().parse(excel_file)
        assert result.page_count == 1

    def test_parse_invalid_file_raises(self, tmp_path):
        bad = tmp_path / "bad.xlsx"
        bad.write_text("not an excel file")
        with pytest.raises(ParseFailureError):
            ExcelParser().parse(bad)


class TestCSVParser:
    def test_format_and_extensions(self):
        p = CSVParser()
        assert p.supported_format == FileFormat.CSV
        assert "csv" in p.supported_extensions

    def test_parse_returns_sheet_keyed_by_filename(self, csv_file):
        result = CSVParser().parse(csv_file)
        assert result.format == FileFormat.CSV
        assert "sample" in result.sheets
        assert result.sheets["sample"][0]["city"] == "Chennai"

    def test_parse_metadata_has_column_info(self, csv_file):
        result = CSVParser().parse(csv_file)
        assert result.metadata["row_count"] == 1
        assert "city" in result.metadata["columns"]


class TestWordParser:
    def test_format_and_extensions(self):
        p = WordParser()
        assert p.supported_format == FileFormat.WORD
        assert "docx" in p.supported_extensions

    def test_parse_returns_text_content(self, word_file):
        result = WordParser().parse(word_file)
        assert result.format == FileFormat.WORD
        assert "This is a test paragraph." in result.text_content

    def test_headings_captured_in_metadata(self, word_file):
        result = WordParser().parse(word_file)
        headings = result.metadata["headings"]
        assert any(h["text"] == "Test Heading" for h in headings)

    def test_parse_invalid_file_raises(self, tmp_path):
        bad = tmp_path / "bad.docx"
        bad.write_text("not a docx")
        with pytest.raises(ParseFailureError):
            WordParser().parse(bad)


class TestPPTParser:
    def test_format_and_extensions(self):
        p = PPTParser()
        assert p.supported_format == FileFormat.PPT
        assert "pptx" in p.supported_extensions

    def test_parse_returns_slides(self, ppt_file):
        result = PPTParser().parse(ppt_file)
        assert result.format == FileFormat.PPT
        assert result.page_count == 1
        assert len(result.slides) == 1
        assert result.slides[0]["slide_number"] == 1
        assert result.slides[0]["title"] == "Test Slide Title"

    def test_flat_text_content_populated(self, ppt_file):
        result = PPTParser().parse(ppt_file)
        assert len(result.text_content) == 1
        assert "Test Slide Title" in result.text_content[0]
