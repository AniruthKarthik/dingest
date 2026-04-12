import hashlib
import tempfile
import io
from pathlib import Path
from fastapi import APIRouter, Depends, UploadFile, File, status
from fastapi.responses import StreamingResponse
from app.models.document import ParsedDocument
from app.services.ingestion_service import IngestionService
from app.utils.file_validator import validate_upload

router = APIRouter(prefix="/ingest", tags=["ingestion"])


def get_ingestion_service() -> IngestionService:
    return IngestionService()


@router.post(
    "/upload",
    response_model=ParsedDocument,
    status_code=status.HTTP_200_OK,
    summary="Upload and parse a document",
)
async def upload_and_parse(
    file: UploadFile = File(...),
    service: IngestionService = Depends(get_ingestion_service),
) -> ParsedDocument:
    content = await file.read()
    validate_upload(file, content)

    suffix = Path(file.filename).suffix
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(content)
        tmp_path = Path(tmp.name)

    try:
        result = service.ingest(tmp_path)
        result.file_name = file.filename
        result.file_size_bytes = len(content)
        result.file_hash = hashlib.shake_256(content).hexdigest(16)

        return result
    finally:
        tmp_path.unlink(missing_ok=True)


@router.get("/formats")
def supported_formats() -> dict:
    from app.ingestion.factory import ParserFactory
    return {"supported_extensions": sorted(ParserFactory.supported_extensions())}


@router.post("/export")
async def export_document(doc: ParsedDocument):
    """Reconstructs a binary file from ParsedDocument JSON (with edits)."""
    output = io.BytesIO()
    media_type = "application/octet-stream"

    if doc.format == "excel" or doc.format == "csv":
        import pandas as pd
        if doc.sheets:
            with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
                for sheet_name, rows in doc.sheets.items():
                    df = pd.DataFrame(rows)
                    df.to_excel(writer, sheet_name=sheet_name, index=False)
            media_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"

    elif doc.format == "word":
        from docx import Document
        new_doc = Document()
        if doc.title:
            new_doc.add_heading(doc.title, 0)
        if doc.text_content:
            for block in doc.text_content:
                new_doc.add_paragraph(block)
        new_doc.save(output)
        media_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

    elif doc.format == "ppt":
        from pptx import Presentation
        prs = Presentation()
        slide_layout = prs.slide_layouts[1] # Title and Content
        if doc.slides:
            for s in doc.slides:
                slide = prs.slides.add_slide(slide_layout)
                if s.get("title"):
                    slide.shapes.title.text = s["title"]
                if s.get("content"):
                    slide.placeholders[1].text = s["content"]
        prs.save(output)
        media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    elif doc.format == "pdf":
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
        c = canvas.Canvas(output, pagesize=letter)
        width, height = letter
        if doc.text_content:
            for page_text in doc.text_content:
                text_obj = c.beginText(50, height - 50)
                text_obj.setFont("Helvetica", 10)
                # Split text into lines to avoid overflow
                for line in page_text.splitlines():
                    text_obj.textLine(line)
                c.drawText(text_obj)
                c.showPage()
        c.save()
        media_type = "application/pdf"

    output.seek(0)
    return StreamingResponse(
        output,
        media_type=media_type,
        headers={"Content-Disposition": f"attachment; filename={doc.file_name}"}
    )
